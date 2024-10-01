import os
import streamlit as st
from groq import Groq
import pdfplumber
import requests
from dotenv import load_dotenv
import subprocess
import io
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document

load_dotenv()

api_key = os.getenv("GROQ_API_KEY")
hf_api_key = os.getenv("HF_API_KEY")
image_api_key = os.getenv("IMAGE_API_KEY")

if not api_key:
    raise ValueError("GROQ_API_KEY environment variable is not set")

# Initialize Groq client with the API key
client = Groq(api_key=api_key)

# Set API URLs and headers with API key from environment variables
HF_API_URL = "https://api-inference.huggingface.co/models/facebook/bart-large-cnn"
HF_API_HEADERS = {"Authorization": f"Bearer {hf_api_key}"}

IMAGE_API_URL = "https://api-inference.huggingface.co/models/gpt153/SAM"
IMAGE_API_HEADERS = {"Authorization": f"Bearer {image_api_key}"}

def summarize_text_hf_api(text, max_chunk_size=2000):
    payload = {"inputs": text[:max_chunk_size]}  # Limiting text chunk size
    response = requests.post(HF_API_URL, headers=HF_API_HEADERS, json=payload)
    
    if response.status_code == 200:
        summary = response.json()
        return summary[0]["summary_text"] if summary and "summary_text" in summary[0] else "No summary available."
    else:
        return f"Error: {response.status_code}, {response.text}"

def extract_text_from_pdf(uploaded_file):
    with pdfplumber.open(uploaded_file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(uploaded_file):
    doc = Document(uploaded_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def convert_doc_to_docx(doc_path):
    docx_path = doc_path.replace(".doc", ".docx")
    subprocess.run(["unoconv", "-f", "docx", doc_path], check=True)
    return docx_path

def extract_text_from_doc(uploaded_file):
    doc_path = os.path.join("/tmp", uploaded_file.name)
    with open(doc_path, "wb") as f:
        f.write(uploaded_file.getvalue())
    
    docx_path = convert_doc_to_docx(doc_path)
    
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    
    return text

def generate_image_from_text(prompt):
    payload = {"inputs": prompt}
    response = requests.post(IMAGE_API_URL, headers=IMAGE_API_HEADERS, json=payload)
    
    if response.status_code == 200:
        return response.content
    else:
        st.error(f"Image generation failed: {response.status_code}, {response.text}")
        return None

def generate_powerpoint(summary):
    prs = Presentation()
    for section, content in summary.items():
        slide_layout = prs.slide_layouts[5]  
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = section
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        image_bytes = generate_image_from_text(content)

        left_textbox = Inches(0.5)
        top_textbox = Inches(1.5)
        width_textbox = Inches(4)
        height_textbox = Inches(5)
        textbox = slide.shapes.add_textbox(left_textbox, top_textbox, width_textbox, height_textbox)
        content_box = textbox.text_frame

        content_box.text = content
        for paragraph in content_box.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
            paragraph.alignment = PP_ALIGN.LEFT

        if image_bytes:
            image = Image.open(io.BytesIO(image_bytes))
            image_path = "generated_image.png"
            image.save(image_path)

            left_image = Inches(5.5)
            top_image = Inches(1.5)
            width_image = Inches(4)
            slide.shapes.add_picture(image_path, left_image, top_image, width=width_image)

    return prs

def generate_text_with_groq_streaming(prompt):
    completion = client.chat.completions.create(
        model="llama-3.1-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        temperature=1,
        max_tokens=7000,
        top_p=1,
        stream=True,
        stop=None,
    )
    full_text = ""
    text_placeholder = st.empty()
    for chunk in completion:
        content_chunk = chunk.choices[0].delta.content or ""
        full_text += content_chunk
        text_placeholder.text(full_text)

def generate_response_with_groq(prompt, context):
    completion = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {"role": "system", "content": f"Research content: {context}"},
            {"role": "user", "content": prompt}
        ],
        temperature=1,
        max_tokens=1024,
        top_p=1,
        stream=False,
        stop=None,
    )
    return completion.choices[0].message.content

st.sidebar.title("Choose Functionality")
option = st.sidebar.radio("Select an option", ["Text Generation", "Text Summarization", "Data Visualization", "Chat with PDF"])

if option == "Text Generation":
    st.title("Question-answering chatbot")
    prompt = st.text_input("Enter any prompt:")
    if st.button("Generate Text"):
        if prompt:
            generate_text_with_groq_streaming(prompt)
        else:
            st.warning("Please enter a prompt.")

elif option == "Text Summarization":
    st.title("Upload document for Summarization")
    uploaded_file = st.file_uploader("Upload your Document", type=["pdf", "docx", "doc"])
    
    if uploaded_file:
        if uploaded_file.type == "application/pdf":
            st.write("Extracting text from PDF and summarizing the document...")
            paper_text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            st.write("Extracting text from DOCX and summarizing the document...")
            paper_text = extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "application/msword":
            st.write("Extracting text from DOC and summarizing the document...")
            paper_text = extract_text_from_doc(uploaded_file)
        else:
            st.error("Unsupported file type. Please upload a PDF, DOCX, or DOC file.")
            paper_text = ""
        
        if paper_text:
            st.write("Summarizing the document...")
            summary = summarize_text_hf_api(paper_text)
            st.write("**Summary:**")
            st.write(summary)

elif option == "Data Visualization":
    st.title("Upload Dataset for Visualization")
    uploaded_file = st.file_uploader("Upload your CSV file", type=["csv"])
    if uploaded_file:
        data = pd.read_csv(uploaded_file)
        st.write("Data preview:")
        st.write(data.head())
        
        chart_type = st.selectbox("Select chart type", ["Bar", "Line", "Pie"])
        
        if st.button("Generate Chart"):
            if chart_type == "Bar":
                fig, ax = plt.subplots()
                data.plot(kind='bar', ax=ax, color='skyblue')
                plt.title("Bar Chart")
                plt.xlabel("X-axis")
                plt.ylabel("Y-axis")
            elif chart_type == "Line":
                fig, ax = plt.subplots()
                data.plot(kind='line', ax=ax, color='skyblue')
                plt.title("Line Chart")
                plt.xlabel("X-axis")
                plt.ylabel("Y-axis")
            elif chart_type == "Pie":
                fig, ax = plt.subplots()
                data.plot(kind='pie', ax=ax, y=data.columns[0], autopct='%1.1f%%')
                plt.title("Pie Chart")
                
            chart_image_path = "chart.png"
            plt.savefig(chart_image_path)
            plt.close()
            
            st.image(chart_image_path, caption="Generated Chart")
            
            with open(chart_image_path, "rb") as file:
                st.download_button("Download Chart", file, file_name="chart.png")

elif option == "Chat with PDF":
    st.title("Chat with any document of your choice")

    if 'conversation' not in st.session_state:
        st.session_state['conversation'] = []

    uploaded_file = st.file_uploader("Upload your document", type=["pdf", "docx", "doc"])

    if uploaded_file:
        if uploaded_file.type == "application/pdf":
            st.session_state['research_content'] = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            st.session_state['research_content'] = extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "application/msword":
            st.session_state['research_content'] = extract_text_from_doc(uploaded_file)
        else:
            st.error("Unsupported file type. Please upload a PDF, DOCX, or DOC file.")
        
        if 'research_content' in st.session_state:
            st.write("**Preview of your uploaded document:**")
            preview_length = 100000 
            st.text_area("File Preview", st.session_state['research_content'][:preview_length], height=300)
            st.write("You can now ask questions related to your research paper.")

    st.subheader("Ask questions about your document")

    user_input = st.text_input("Enter your question related to the document:")

    if st.button("Ask"):
        if 'research_content' in st.session_state and user_input:
            response = generate_response_with_groq(user_input, st.session_state['research_content'])
            
            st.session_state['conversation'].append({"user": user_input, "assistant": response})

    if st.session_state['conversation']:
        st.write("Conversation History:")
        for chat in st.session_state['conversation']:
            st.markdown(f"**You:** {chat['user']}")
            st.markdown(f"**Assistant:** {chat['assistant']}")

    if st.button("Clear Conversation History"):
        st.session_state['conversation'] = []
        st.write("Conversation history cleared.")
